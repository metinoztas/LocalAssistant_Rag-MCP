"""
Document Loader

Görevleri:
- Desteklenen belge türlerini okumak
- Dosya uzantısına göre doğru okuyucuyu seçmek
- Belge içeriğini metin olarak döndürmek
"""

from pathlib import Path

import fitz
import markdown

from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation


def _load_txt(file_path: str) -> str:
    """
    TXT dosyasını okur.
    """

    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()


def _load_pdf(file_path: str) -> str:
    """
    PDF dosyasını okur.
    """


    pages = []

    with fitz.open(file_path) as document:

        for page in document:

            text = page.get_text("text")

            if text.strip():
                pages.append(text)

    return "\n\n".join(pages)


def _load_docx(file_path: str) -> str:
    """
    DOCX dosyasını okur.
    """

    document = Document(file_path)

    paragraphs = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    return "\n\n".join(paragraphs)


def _load_markdown(file_path: str) -> str:
    """
    Markdown dosyasını okur.
    """

    with open(file_path, "r", encoding="utf-8") as file:
        markdown_text = file.read()

    html = markdown.markdown(markdown_text)

    soup = BeautifulSoup(html, "html.parser")

    return soup.get_text(separator="\n")


def _load_html(file_path: str) -> str:
    """
    HTML dosyasını okur.
    """

    with open(file_path, "r", encoding="utf-8") as file:
        html = file.read()

    soup = BeautifulSoup(html, "html.parser")

    return soup.get_text(separator="\n")


def _load_pptx(file_path: str) -> str:
    """
    PowerPoint dosyasını okur.
    """

    presentation = Presentation(file_path)

    slides = []

    for slide in presentation.slides:

        texts = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    texts.append(text)

        if texts:
            slides.append("\n".join(texts))

    return "\n\n".join(slides)


DOCUMENT_LOADERS = {
    ".txt": _load_txt,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".md": _load_markdown,
    ".markdown": _load_markdown,
    ".html": _load_html,
    ".htm": _load_html,
    ".pptx": _load_pptx,
}


def load_document(file_path: str) -> str:
    """
    Dosya uzantısına göre uygun belge okuyucusunu çalıştırır.
    """

    extension = Path(file_path).suffix.lower()

    loader = DOCUMENT_LOADERS.get(extension)

    if loader is None:

        supported = ", ".join(sorted(DOCUMENT_LOADERS.keys()))

        raise ValueError(
            f"Desteklenmeyen dosya türü: {extension}\n"
            f"Desteklenen türler: {supported}"
        )

    return loader(file_path)